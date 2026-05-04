#!/usr/bin/env python3
"""
KU Automation Portfolio Presentation Generator
Creates a professional, optimized PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Brand colors
PRIMARY = RGBColor(79, 70, 229)  # Indigo #4F46E5
PRIMARY_DARK = RGBColor(49, 46, 129)  # Dark indigo
ACCENT = RGBColor(139, 92, 246)  # Purple
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(17, 24, 39)  # Gray 900
GRAY = RGBColor(107, 114, 128)  # Gray 500
LIGHT_BG = RGBColor(249, 250, 251)  # Gray 50

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add a text box with specified properties"""
    if color is None:
        color = DARK
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: COVER ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, PRIMARY_DARK)
    bg.line.fill.background()
    
    # Logo box
    logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5))
    set_shape_fill(logo_box, WHITE)
    logo_box.line.fill.background()
    
    add_text_box(slide, 5.9, 1.85, 1.5, 0.8, "KU", 36, True, PRIMARY, PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide, 0, 3.3, 13.333, 1, "KU Automation", 54, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, 0, 4.2, 13.333, 0.6, "AI-Powered Engineering Solutions", 28, False, RGBColor(200, 200, 220), PP_ALIGN.CENTER)
    
    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(5), Inches(2.5), Inches(0.05))
    set_shape_fill(line, ACCENT)
    line.line.fill.background()
    
    add_text_box(slide, 0, 5.3, 13.333, 0.5, "Corporate Portfolio & Capabilities", 18, False, RGBColor(180, 180, 200), PP_ALIGN.CENTER)
    add_text_box(slide, 0, 6.5, 13.333, 0.4, "May 2026", 14, False, RGBColor(140, 140, 160), PP_ALIGN.CENTER)
    
    # ========== SLIDE 2: ABOUT US ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_text_box(slide, 0.5, 0.4, 3, 0.3, "01 — ABOUT US", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 6, 0.8, "Who We Are", 40, True, DARK)
    
    # Founder box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.8))
    set_shape_fill(box1, LIGHT_BG)
    box1.line.fill.background()
    
    add_text_box(slide, 0.8, 1.9, 5.2, 0.4, "Kingsley Uzowulu", 20, True, DARK)
    add_text_box(slide, 0.8, 2.3, 5.2, 0.3, "Founder & Lead Engineer", 14, False, PRIMARY)
    add_text_box(slide, 0.8, 2.7, 5.2, 1.5, "Chartered Engineer (CEng MIMechE) with 21+ years of experience in oil & gas, EPC projects, and manufacturing across the UK, Europe, and Middle East.", 12, False, GRAY)
    
    # Tags
    tag1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(1.7), Inches(0.35))
    set_shape_fill(tag1, RGBColor(224, 231, 255))
    tag1.line.fill.background()
    add_text_box(slide, 0.85, 4.03, 1.6, 0.3, "Chartered Engineer", 9, False, PRIMARY)
    
    tag2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), Inches(4.0), Inches(1.0), Inches(0.35))
    set_shape_fill(tag2, RGBColor(224, 231, 255))
    tag2.line.fill.background()
    add_text_box(slide, 2.65, 4.03, 0.9, 0.3, "MIMechE", 9, False, PRIMARY)
    
    tag3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(4.0), Inches(1.4), Inches(0.35))
    set_shape_fill(tag3, RGBColor(224, 231, 255))
    tag3.line.fill.background()
    add_text_box(slide, 3.75, 4.03, 1.3, 0.3, "AI/ML Specialist", 9, False, PRIMARY)
    
    # Right column
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.8), Inches(1.3))
    set_shape_fill(box2, LIGHT_BG)
    box2.line.fill.background()
    add_text_box(slide, 7.1, 1.85, 5.2, 0.35, "🎯 Our Mission", 14, True, DARK)
    add_text_box(slide, 7.1, 2.25, 5.2, 0.7, "Transform engineering productivity through AI automation that actually works in real-world project environments.", 11, False, GRAY)
    
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.1), Inches(5.8), Inches(2.0))
    set_shape_fill(box3, LIGHT_BG)
    box3.line.fill.background()
    add_text_box(slide, 7.1, 3.25, 5.2, 0.35, "🏭 Industries Served", 14, True, DARK)
    add_text_box(slide, 7.1, 3.65, 5.2, 1.3, "• Oil & Gas (Upstream, Midstream, Downstream)\n• EPC Contractors & Engineering Consultancies\n• Chemical, Power Generation & Manufacturing", 10, False, GRAY)
    
    # Quote box
    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.6), Inches(5.8), Inches(1.0))
    set_shape_fill(quote_box, RGBColor(238, 242, 255))
    quote_box.line.fill.background()
    
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.6), Inches(0.08), Inches(1.0))
    set_shape_fill(accent_line, PRIMARY)
    accent_line.line.fill.background()
    
    add_text_box(slide, 0.8, 4.75, 5.2, 0.7, '"We bridge the gap between cutting-edge AI and practical engineering workflows."', 12, False, PRIMARY_DARK)
    
    add_text_box(slide, 6.8, 5.3, 5.8, 0.35, "🌍 Geographic Coverage", 14, True, DARK)
    add_text_box(slide, 6.8, 5.7, 5.8, 0.35, "UK • Europe • Middle East • Global Remote Support", 12, False, GRAY)
    
    # ========== SLIDE 3: THE PROBLEM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, RGBColor(17, 24, 39))
    bg.line.fill.background()
    
    add_text_box(slide, 0.5, 0.4, 4, 0.3, "02 — THE CHALLENGE", 11, True, RGBColor(129, 140, 248))
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "The Problem We Solve", 40, True, WHITE)
    
    stats = [
        ("60%", "of senior engineer time on document admin", RGBColor(248, 113, 113)),
        ("£50B", "annual cost of unplanned downtime (UK)", RGBColor(251, 191, 36)),
        ("4-6 hrs", "to manually review a single tech spec", RGBColor(251, 146, 60))
    ]
    
    for i, (num, desc, color) in enumerate(stats):
        x = 0.5 + i * 4.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(3.8), Inches(1.8))
        set_shape_fill(box, RGBColor(30, 41, 59))
        box.line.color.rgb = RGBColor(55, 65, 81)
        
        add_text_box(slide, x + 0.3, 2.0, 3.2, 0.7, num, 42, True, color)
        add_text_box(slide, x + 0.3, 2.7, 3.2, 0.7, desc, 11, False, RGBColor(200, 200, 210))
    
    pain_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(12.333), Inches(2.8))
    set_shape_fill(pain_box, RGBColor(30, 41, 59))
    pain_box.line.color.rgb = RGBColor(55, 65, 81)
    
    add_text_box(slide, 0.8, 4.2, 5, 0.4, "Common Pain Points We Address", 16, True, WHITE)
    
    add_text_box(slide, 0.8, 4.7, 5.5, 1.6, "✗ Manual document review bottlenecks\n✗ Inconsistent data extraction from drawings\n✗ Delayed RFQ responses losing contracts", 12, False, RGBColor(200, 200, 210))
    add_text_box(slide, 6.5, 4.7, 5.5, 1.6, "✗ Knowledge trapped in experts' heads\n✗ Contractor deliverable review backlog\n✗ Compliance documentation burden", 12, False, RGBColor(200, 200, 210))
    
    # ========== SLIDE 4: SOLUTIONS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_text_box(slide, 0.5, 0.4, 4, 0.3, "03 — OUR SOLUTIONS", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "AI Solutions Portfolio", 40, True, DARK)
    add_text_box(slide, 0.5, 1.35, 8, 0.4, "End-to-end AI automation for engineering workflows", 14, False, GRAY)
    
    solutions = [
        ("📊", "P&ID Data Extraction", "Extract equipment, valves, instruments", "70% time reduction", RGBColor(59, 130, 246)),
        ("📄", "Datasheet Processing", "Parse datasheets into structured data", "85% accuracy gain", RGBColor(16, 185, 129)),
        ("⚡", "RFQ/Tender Analysis", "AI-powered RFQ analysis", "3x faster responses", RGBColor(168, 85, 247)),
        ("🔍", "Document Review", "6-layer assessment + code refs", "4 hrs → 4 mins", RGBColor(6, 182, 212)),
        ("📈", "Stress Analysis Review", "CAESAR II / AutoPIPE validation", "Catch software misses", RGBColor(139, 92, 246)),
        ("🤖", "Engineering Chatbots", "Knowledge base on your standards", "24/7 answers", RGBColor(245, 158, 11))
    ]
    
    for i, (icon, title, desc, result, color) in enumerate(solutions):
        row = i // 3
        col = i % 3
        x = 0.5 + col * 4.2
        y = 1.9 + row * 2.5
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.9), Inches(2.2))
        set_shape_fill(box, RGBColor(249, 250, 251))
        box.line.fill.background()
        
        add_text_box(slide, x + 0.2, y + 0.15, 0.5, 0.5, icon, 24, False, DARK)
        add_text_box(slide, x + 0.2, y + 0.65, 3.5, 0.4, title, 14, True, DARK)
        add_text_box(slide, x + 0.2, y + 1.05, 3.5, 0.55, desc, 10, False, GRAY)
        add_text_box(slide, x + 0.2, y + 1.6, 3.5, 0.3, result, 11, True, color)
    
    # ========== SLIDE 5: LIVE DEMOS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    add_text_box(slide, 0.5, 0.4, 4, 0.3, "04 — LIVE DEMONSTRATIONS", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Production-Ready Tools", 40, True, DARK)
    add_text_box(slide, 0.5, 1.35, 8, 0.4, "Try our tools live — no signup required", 14, False, GRAY)
    
    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(6), Inches(4.8))
    set_shape_fill(demo_box, WHITE)
    demo_box.line.fill.background()
    
    add_text_box(slide, 0.8, 2.1, 5.4, 0.3, "🟢 Live at www.ku-automation.com/demos", 11, False, RGBColor(22, 163, 74))
    add_text_box(slide, 0.8, 2.5, 5.4, 0.4, "Available Demos", 18, True, DARK)
    
    demos = "📊 P&ID Parser v2.0\n📄 Datasheet Parser v2.0\n⚡ RFQ Analyzer v2.0\n🔍 Piping Document Review v2.0\n📈 Stress Analysis Review (NEW)"
    add_text_box(slide, 0.8, 3.0, 5.4, 2.5, demos, 12, False, DARK)
    
    cap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(6), Inches(4.8))
    set_shape_fill(cap_box, WHITE)
    cap_box.line.fill.background()
    
    add_text_box(slide, 7.1, 2.1, 5.4, 0.4, "Demo Capabilities", 18, True, DARK)
    add_text_box(slide, 7.1, 2.6, 5.4, 0.3, "Input Formats", 12, True, DARK)
    add_text_box(slide, 7.1, 2.9, 5.4, 0.3, "PDF • JPG/PNG • Multi-page documents", 11, False, GRAY)
    add_text_box(slide, 7.1, 3.4, 5.4, 0.3, "Output Formats", 12, True, DARK)
    add_text_box(slide, 7.1, 3.7, 5.4, 0.3, "Excel (.xlsx) • JSON • REST API", 11, False, GRAY)
    add_text_box(slide, 7.1, 4.2, 5.4, 0.3, "Processing Speed", 12, True, DARK)
    
    speed_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(4.5), Inches(5.4), Inches(1.2))
    set_shape_fill(speed_box, RGBColor(238, 242, 255))
    speed_box.line.fill.background()
    add_text_box(slide, 7.3, 4.65, 5, 0.5, "30-90 seconds", 28, True, PRIMARY)
    add_text_box(slide, 7.3, 5.15, 5, 0.4, "per document (depending on complexity)", 11, False, GRAY)
    
    # ========== SLIDE 6: DEPLOYMENT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_text_box(slide, 0.5, 0.4, 4, 0.3, "05 — DEPLOYMENT", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Flexible Deployment Options", 40, True, DARK)
    
    deployments = [
        ("☁️", "Cloud SaaS", "Fully managed on secure cloud", "No IT infrastructure required\nAutomatic updates & maintenance\nWeb browser access\nAPI integration available\n99.9% uptime SLA", RGBColor(59, 130, 246), RGBColor(219, 234, 254), True),
        ("🔐", "Private Cloud", "Dedicated instance on your cloud", "Dedicated resources\nData stays in your region\nVPC integration\nCustom security policies\nSSO/SAML support", RGBColor(168, 85, 247), RGBColor(243, 232, 255), False),
        ("🏢", "On-Premise", "Full deployment in your data center", "Complete data control\nAir-gapped option available\nIntegrate with existing systems\nCustom hardware specs\nPerpetual license option", RGBColor(107, 114, 128), RGBColor(243, 244, 246), False)
    ]
    
    for i, (icon, title, desc, features, color, bg_color, popular) in enumerate(deployments):
        x = 0.5 + i * 4.2
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.9), Inches(4.5))
        set_shape_fill(box, bg_color)
        box.line.color.rgb = color
        
        add_text_box(slide, x + 0.3, 1.8, 0.6, 0.6, icon, 28, False, DARK)
        add_text_box(slide, x + 0.3, 2.4, 3.3, 0.4, title, 18, True, DARK)
        add_text_box(slide, x + 0.3, 2.8, 3.3, 0.4, desc, 10, False, GRAY)
        
        # Features as bullet points
        feature_text = "\n".join([f"✓ {f}" for f in features.split("\n")])
        add_text_box(slide, x + 0.3, 3.3, 3.3, 2.5, feature_text, 10, False, DARK)
        
        if popular:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1), Inches(1.4), Inches(1.8), Inches(0.35))
            set_shape_fill(badge, color)
            badge.line.fill.background()
            add_text_box(slide, x + 1.05, 1.43, 1.7, 0.3, "MOST POPULAR", 8, True, WHITE, PP_ALIGN.CENTER)
    
    int_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.9))
    set_shape_fill(int_box, RGBColor(238, 242, 255))
    int_box.line.fill.background()
    add_text_box(slide, 0.8, 6.45, 12, 0.5, "🔌 Integrations: SharePoint • SAP • AVEVA/Hexagon • REST API • Custom integrations", 11, False, DARK)
    
    # ========== SLIDE 7: SECURITY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, RGBColor(17, 24, 39))
    bg.line.fill.background()
    
    add_text_box(slide, 0.5, 0.4, 5, 0.3, "06 — SECURITY & COMPLIANCE", 11, True, RGBColor(129, 140, 248))
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Enterprise-Grade Security", 40, True, WHITE)
    
    dp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.9), Inches(3.2))
    set_shape_fill(dp_box, RGBColor(30, 41, 59))
    dp_box.line.color.rgb = RGBColor(55, 65, 81)
    
    add_text_box(slide, 0.8, 1.8, 5.3, 0.4, "🔒 Data Protection", 16, True, WHITE)
    add_text_box(slide, 0.8, 2.3, 5.3, 2.5, "✓ End-to-End Encryption (AES-256, TLS 1.3)\n✓ Data Isolation (Tenant-isolated databases)\n✓ No Data Retention (Configurable deletion)\n✓ Data Residency (UK, EU, or regional)", 12, False, RGBColor(180, 220, 180))
    
    ac_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.9), Inches(3.2))
    set_shape_fill(ac_box, RGBColor(30, 41, 59))
    ac_box.line.color.rgb = RGBColor(55, 65, 81)
    
    add_text_box(slide, 7.1, 1.8, 5.3, 0.4, "👤 Access Control", 16, True, WHITE)
    add_text_box(slide, 7.1, 2.3, 5.3, 2.5, "✓ Single Sign-On (SAML, OAuth 2.0, Azure AD)\n✓ Role-Based Access (Granular permissions)\n✓ Multi-Factor Authentication\n✓ Audit Logging (Complete activity trails)", 12, False, RGBColor(180, 220, 180))
    
    comp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(12.333), Inches(1.6))
    set_shape_fill(comp_box, RGBColor(30, 41, 59))
    comp_box.line.color.rgb = RGBColor(55, 65, 81)
    
    add_text_box(slide, 0.8, 5.3, 12, 0.4, "Compliance & Certifications", 14, True, WHITE)
    add_text_box(slide, 0.8, 5.8, 12, 0.6, "🇬🇧 UK GDPR        🔐 ISO 27001        ☁️ SOC 2 Type II        🏭 Industry Standards", 14, False, WHITE)
    
    # ========== SLIDE 8: SUPPORT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_text_box(slide, 0.5, 0.4, 5, 0.3, "07 — SUPPORT & MAINTENANCE", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Ongoing Support Model", 40, True, DARK)
    
    support_tiers = [
        ("Standard", "Included with all plans", "✓ Email support (24hr response)\n✓ Documentation & knowledge base\n✓ Quarterly software updates\n✓ Bug fixes & security patches", LIGHT_BG, DARK),
        ("Priority", "Recommended for production", "✓ 4-hour response time\n✓ Phone & video call support\n✓ Dedicated account manager\n✓ Monthly check-in calls\n✓ Priority feature requests", RGBColor(238, 242, 255), DARK),
        ("Enterprise", "Mission-critical deployments", "✓ 1-hour response (24/7)\n✓ On-site support available\n✓ Custom SLA agreements\n✓ Dedicated engineering team\n✓ Quarterly business reviews", RGBColor(17, 24, 39), WHITE)
    ]
    
    for i, (name, desc, features, bg_color, text_color) in enumerate(support_tiers):
        x = 0.5 + i * 4.2
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.9), Inches(4.0))
        set_shape_fill(box, bg_color)
        if i == 1:
            box.line.color.rgb = PRIMARY
        else:
            box.line.fill.background()
        
        add_text_box(slide, x + 0.3, 1.8, 3.3, 0.4, name, 18, True, text_color)
        add_text_box(slide, x + 0.3, 2.2, 3.3, 0.3, desc, 10, False, GRAY if i < 2 else RGBColor(180, 180, 200))
        add_text_box(slide, x + 0.3, 2.7, 3.3, 2.5, features, 10, False, text_color)
    
    cycle_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
    set_shape_fill(cycle_box, RGBColor(238, 242, 255))
    cycle_box.line.fill.background()
    add_text_box(slide, 0.8, 6.0, 12, 0.8, "🔄 Continuous Improvement:  1. Monitor  →  2. Optimize  →  3. Update  →  4. Train", 14, False, DARK)
    
    # ========== SLIDE 9: ENGAGEMENT MODEL ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    add_text_box(slide, 0.5, 0.4, 4, 0.3, "08 — ENGAGEMENT MODEL", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "How We Work Together", 40, True, DARK)
    
    phases = [
        ("1", "Week 1-2", "Discovery & Assessment", "Understand workflows, review samples, IT assessment"),
        ("2", "Week 3-4", "Proof of Concept", "Configure with your data, pilot, measure results"),
        ("3", "Week 5-8", "Production Deployment", "Full deployment, integration, training, go-live"),
        ("✓", "Ongoing", "Continuous Support", "Check-ins, updates, scale to more use cases")
    ]
    
    for i, (num, time, title, desc) in enumerate(phases):
        y = 1.5 + i * 1.4
        
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.7), Inches(0.7))
        if i < 3:
            set_shape_fill(circle, PRIMARY)
        else:
            set_shape_fill(circle, RGBColor(22, 163, 74))
        circle.line.fill.background()
        add_text_box(slide, 0.78, y + 0.15, 0.55, 0.4, num, 18, True, WHITE, PP_ALIGN.CENTER)
        
        content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(y - 0.05), Inches(10.5), Inches(1.0))
        set_shape_fill(content_box, WHITE)
        content_box.line.fill.background()
        
        add_text_box(slide, 2.0, y + 0.05, 2, 0.3, time, 11, True, PRIMARY)
        add_text_box(slide, 2.0, y + 0.35, 4, 0.35, title, 14, True, DARK)
        add_text_box(slide, 6.3, y + 0.2, 5.5, 0.6, desc, 11, False, GRAY)
    
    # ========== SLIDE 10: RESULTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_text_box(slide, 0.5, 0.4, 3, 0.3, "09 — RESULTS", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Typical Results", 40, True, DARK)
    
    results = [
        ("70%", "Reduction in document\nprocessing time", RGBColor(59, 130, 246)),
        ("85%", "Improvement in\ndata accuracy", RGBColor(16, 185, 129)),
        ("3x", "Faster RFQ\nresponse times", RGBColor(168, 85, 247)),
        ("<4wk", "Typical ROI\npayback period", RGBColor(245, 158, 11))
    ]
    
    for i, (num, desc, color) in enumerate(results):
        x = 0.5 + i * 3.2
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.0), Inches(1.7))
        set_shape_fill(box, color)
        box.line.fill.background()
        
        add_text_box(slide, x + 0.2, 1.65, 2.6, 0.6, num, 36, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 2.35, 2.6, 0.7, desc, 10, False, RGBColor(240, 240, 255), PP_ALIGN.CENTER)
    
    compare_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.5), Inches(12.333), Inches(3.5))
    set_shape_fill(compare_box, LIGHT_BG)
    compare_box.line.fill.background()
    
    add_text_box(slide, 0.8, 3.7, 8, 0.4, "Sample Use Case: Document Review Automation", 16, True, DARK)
    
    add_text_box(slide, 0.8, 4.2, 5.5, 0.35, "❌ Before KU Automation", 14, True, RGBColor(220, 38, 38))
    add_text_box(slide, 0.8, 4.6, 5.5, 2.0, "• 50 isometrics = 25-50 engineer-hours\n• Inconsistent comment quality\n• 60-80% missing code references\n• Contractor disputes", 11, False, GRAY)
    
    add_text_box(slide, 6.8, 4.2, 5.5, 0.35, "✓ After KU Automation", 14, True, RGBColor(22, 163, 74))
    add_text_box(slide, 6.8, 4.6, 5.5, 2.0, "• 50 isometrics = 4-5 hours (90% reduction)\n• Standardized comments across all docs\n• 95%+ include ASME B31.3 references\n• Clear severity ratings (HOLD/COMMENT)", 11, False, GRAY)
    
    # ========== SLIDE 11: PRICING ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    
    add_text_box(slide, 0.5, 0.4, 3, 0.3, "10 — INVESTMENT", 11, True, PRIMARY)
    add_text_box(slide, 0.5, 0.7, 8, 0.8, "Flexible Pricing Models", 40, True, DARK)
    
    pricing = [
        ("Pilot Project", "Prove the value", "£5K - £15K", "✓ 4-week implementation\n✓ Single use case focus\n✓ Up to 10 users\n✓ Standard support\n✓ ROI measurement", WHITE, DARK, False),
        ("Professional", "Full production", "£25K - £75K/yr", "✓ 8-week implementation\n✓ Multiple use cases\n✓ Up to 50 users\n✓ Priority support\n✓ Integrations included", PRIMARY, WHITE, True),
        ("Enterprise", "Full transformation", "Custom", "✓ Dedicated team\n✓ Unlimited use cases\n✓ Unlimited users\n✓ 24/7 support\n✓ On-premise option", RGBColor(17, 24, 39), WHITE, False)
    ]
    
    for i, (name, desc, price, features, bg_color, text_color, popular) in enumerate(pricing):
        x = 0.5 + i * 4.2
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.9), Inches(4.8))
        set_shape_fill(box, bg_color)
        box.line.fill.background()
        
        if popular:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.9), Inches(1.3), Inches(2.1), Inches(0.35))
            set_shape_fill(badge, RGBColor(251, 191, 36))
            badge.line.fill.background()
            add_text_box(slide, x + 0.95, 1.33, 2, 0.3, "MOST POPULAR", 9, True, RGBColor(120, 53, 15), PP_ALIGN.CENTER)
        
        add_text_box(slide, x + 0.3, 1.75, 3.3, 0.4, name, 18, True, text_color)
        subtitle_color = GRAY if i == 0 else RGBColor(180, 180, 200)
        add_text_box(slide, x + 0.3, 2.15, 3.3, 0.3, desc, 11, False, subtitle_color)
        add_text_box(slide, x + 0.3, 2.6, 3.3, 0.5, price, 24, True, text_color)
        add_text_box(slide, x + 0.3, 3.2, 3.3, 2.8, features, 10, False, text_color)
    
    add_text_box(slide, 0, 6.5, 13.333, 0.4, "All pricing is indicative. Final quotes based on scope, complexity, and deployment requirements.", 10, False, GRAY, PP_ALIGN.CENTER)
    
    # ========== SLIDE 12: NEXT STEPS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, PRIMARY_DARK)
    bg.line.fill.background()
    
    add_text_box(slide, 0, 0.5, 13.333, 0.3, "LET'S GET STARTED", 11, True, RGBColor(165, 180, 252), PP_ALIGN.CENTER)
    add_text_box(slide, 0, 0.9, 13.333, 0.8, "Next Steps", 44, True, WHITE, PP_ALIGN.CENTER)
    
    steps = [
        ("1", "Discovery Call", "30-minute call to\nunderstand your needs"),
        ("2", "Custom Demo", "See our tools with\nyour actual documents"),
        ("3", "Proposal", "Detailed scope, timeline\n& investment")
    ]
    
    for i, (num, title, desc) in enumerate(steps):
        x = 1.5 + i * 3.8
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.3), Inches(2.2))
        set_shape_fill(box, RGBColor(55, 48, 163))
        box.line.fill.background()
        
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.15), Inches(2.2), Inches(1), Inches(1))
        set_shape_fill(num_circle, RGBColor(79, 70, 229))
        num_circle.line.fill.background()
        add_text_box(slide, x + 1.2, y + 0.3, 0.9, 0.6, num, 28, True, WHITE, PP_ALIGN.CENTER)
        
        add_text_box(slide, x + 0.2, 3.35, 2.9, 0.35, title, 14, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.7, 2.9, 0.6, desc, 10, False, RGBColor(200, 200, 220), PP_ALIGN.CENTER)
    
    # Contact box
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.5), Inches(6.333), Inches(2.2))
    set_shape_fill(contact_box, WHITE)
    contact_box.line.fill.background()
    
    add_text_box(slide, 3.7, 4.7, 5.9, 0.4, "Contact", 20, True, DARK, PP_ALIGN.CENTER)
    add_text_box(slide, 4.0, 5.2, 5.5, 1.3, "👤  Kingsley Uzowulu\n📧  kingsley.uzowulu@ku-automation.com\n🌐  www.ku-automation.com\n📅  calendly.com/kingsley-uzowulu/30min", 11, False, DARK)
    
    add_text_box(slide, 0, 6.9, 13.333, 0.3, "Thank you for your time and consideration.", 11, False, RGBColor(160, 160, 180), PP_ALIGN.CENTER)
    
    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'KU_Automation_Portfolio_Presentation.pptx')
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
